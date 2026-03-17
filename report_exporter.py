"""
Export financial statements and charts to PowerPoint and PDF formats.
"""

import io
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from fpdf import FPDF


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _fmt(value: float) -> str:
    if value < 0:
        return f"(${abs(value):,.2f})"
    return f"${value:,.2f}"


_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_DARK = RGBColor(0x33, 0x33, 0x33)
_PURPLE = RGBColor(0x66, 0x7E, 0xEA)
_GREEN = RGBColor(0x11, 0x99, 0x8E)
_RED = RGBColor(0xEB, 0x33, 0x49)
_LIGHT_BG = RGBColor(0xF0, 0xF0, 0xF8)
_HEADER_BG = RGBColor(0x44, 0x52, 0x7A)


def _set_cell(cell, text, size=10, bold=False, align=PP_ALIGN.LEFT,
              font_color=_DARK, fill_color=None):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = font_color
    run.font.name = "Calibri"
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color


def _add_title_slide(prs, org_name, subtitle="Financial Report"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x2D, 0x32, 0x56)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = org_name
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = _WHITE
    run.font.name = "Calibri"

    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.6), Inches(8), Inches(0.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.size = Pt(20)
    run2.font.color.rgb = RGBColor(0xAA, 0xB0, 0xD0)
    run2.font.name = "Calibri"

    return slide


def _add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x3A, 0x3F, 0x63)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = _WHITE
    run.font.name = "Calibri"
    return slide


def _add_table_slide(prs, title, headers, rows, col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = _DARK
    run.font.name = "Calibri"

    n_rows = len(rows) + 1
    n_cols = len(headers)
    top = Inches(1.1)
    left = Inches(0.5)
    width = Inches(9)
    row_height = min(Inches(0.35), Inches(5.5) / max(n_rows, 1))
    height = row_height * n_rows

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Emu(int(w * 914400))

    for i, h in enumerate(headers):
        align = PP_ALIGN.RIGHT if i > 0 else PP_ALIGN.LEFT
        _set_cell(tbl.cell(0, i), h, size=9, bold=True, align=align,
                  font_color=_WHITE, fill_color=_HEADER_BG)

    for r_idx, row_data in enumerate(rows):
        is_total = any(kw in str(row_data[0]).upper()
                       for kw in ["TOTAL", "CHANGE IN", "ENDING", "NET CHANGE"])
        bg = _LIGHT_BG if r_idx % 2 == 0 else None
        if is_total:
            bg = RGBColor(0xE8, 0xEA, 0xF6)

        for c_idx, val in enumerate(row_data):
            align = PP_ALIGN.RIGHT if c_idx > 0 else PP_ALIGN.LEFT
            _set_cell(tbl.cell(r_idx + 1, c_idx), str(val), size=9,
                      bold=is_total, align=align, fill_color=bg)

    return slide


def _add_chart_slide(prs, title, chart_image_bytes):
    """Add a slide with an embedded chart image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = _DARK
    run.font.name = "Calibri"

    img_stream = io.BytesIO(chart_image_bytes)
    slide.shapes.add_picture(img_stream, Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.2))
    return slide


def _try_chart_image(fig, width=1600, height=900):
    """Try to render a plotly figure to PNG bytes. Returns None on failure."""
    try:
        return fig.to_image(format="png", width=width, height=height, scale=2)
    except Exception:
        return None


# ---------------------------------------------------------------------------
# Chart builders (reusable for both pptx and pdf)
# ---------------------------------------------------------------------------

def _build_revenue_pie(df):
    import plotly.express as px
    non_transfer = df[df["Category"] != "Internal Account Transfer"]
    rev = non_transfer[non_transfer["Amount"] > 0].groupby("Category")["Amount"].sum().reset_index()
    if len(rev) == 0:
        return None
    fig = px.pie(rev, names="Category", values="Amount",
                 color_discrete_sequence=px.colors.qualitative.Set3, hole=0.4)
    fig.update_traces(textposition="inside", textinfo="percent+label")
    fig.update_layout(showlegend=True, margin=dict(t=40, b=40, l=40, r=40),
                      height=700, width=1200, title="Revenue by Category",
                      font=dict(size=14))
    return fig


def _build_expense_pie(df):
    import plotly.express as px
    non_transfer = df[df["Category"] != "Internal Account Transfer"]
    exp = non_transfer[non_transfer["Amount"] < 0].copy()
    exp["AbsAmount"] = exp["Amount"].abs()
    by_cat = exp.groupby("Category")["AbsAmount"].sum().reset_index()
    by_cat.columns = ["Category", "Amount"]
    if len(by_cat) == 0:
        return None
    fig = px.pie(by_cat, names="Category", values="Amount",
                 color_discrete_sequence=px.colors.qualitative.Pastel, hole=0.4)
    fig.update_traces(textposition="inside", textinfo="percent+label")
    fig.update_layout(showlegend=True, margin=dict(t=40, b=40, l=40, r=40),
                      height=700, width=1200, title="Expenses by Category",
                      font=dict(size=14))
    return fig


def _build_monthly_bar(df):
    import plotly.graph_objects as go
    monthly = df.copy()
    monthly["Month"] = monthly["Date"].dt.to_period("M").dt.to_timestamp()
    m_rev = monthly[monthly["Amount"] > 0].groupby("Month")["Amount"].sum()
    m_exp = monthly[monthly["Amount"] < 0].groupby("Month")["Amount"].sum().abs()
    months = sorted(set(m_rev.index) | set(m_exp.index))
    if not months:
        return None
    fig = go.Figure()
    fig.add_trace(go.Bar(
        x=[m.strftime("%b %Y") for m in months],
        y=[m_rev.get(m, 0) for m in months],
        name="Revenue", marker_color="#38ef7d",
    ))
    fig.add_trace(go.Bar(
        x=[m.strftime("%b %Y") for m in months],
        y=[m_exp.get(m, 0) for m in months],
        name="Expenses", marker_color="#f45c43",
    ))
    fig.update_layout(barmode="group", height=700, width=1200,
                      title="Monthly Cash Flow",
                      margin=dict(t=60, b=60, l=60, r=40),
                      font=dict(size=14),
                      legend=dict(orientation="h", yanchor="bottom", y=1.02,
                                  xanchor="right", x=1))
    return fig


# ---------------------------------------------------------------------------
# PowerPoint export
# ---------------------------------------------------------------------------

def export_to_pptx(statements: dict, transactions_df: pd.DataFrame,
                   org_name: str = "Organization") -> bytes:
    """Export financial statements and charts to a PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    act = statements["activities"]
    pos = statements["position"]
    func = statements["functional_expenses"]
    cf = statements["cash_flows"]

    period = act.get("period", "")

    # --- Title slide ---
    _add_title_slide(prs, org_name, period)

    # --- Summary metrics slide ---
    d = act["without_donor_restrictions"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    p = txBox.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Financial Summary"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.name = "Calibri"

    metrics = [
        ("Total Revenue", d["total_revenue"], _GREEN),
        ("Total Expenses", d["total_expenses"], _RED),
        ("Change in Net Assets", d["change_in_net_assets"], _PURPLE),
    ]
    for i, (label, value, color) in enumerate(metrics):
        left = Inches(0.5 + i * 3.1)
        shape = slide.shapes.add_shape(
            1, left, Inches(1.5), Inches(2.8), Inches(2.0)  # MSO_SHAPE.RECTANGLE
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = label
        r1.font.size = Pt(14)
        r1.font.color.rgb = _WHITE
        r1.font.name = "Calibri"

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(12)
        r2 = p2.add_run()
        r2.text = _fmt(value)
        r2.font.size = Pt(28)
        r2.font.bold = True
        r2.font.color.rgb = _WHITE
        r2.font.name = "Calibri"

    # --- Chart slides ---
    charts = [
        ("Revenue by Category", _build_revenue_pie),
        ("Expenses by Category", _build_expense_pie),
        ("Monthly Cash Flow", _build_monthly_bar),
    ]
    for chart_title, builder in charts:
        fig = builder(transactions_df)
        if fig is not None:
            img_bytes = _try_chart_image(fig)
            if img_bytes:
                _add_chart_slide(prs, chart_title, img_bytes)

    # --- Statement of Activities ---
    _add_section_slide(prs, "Statement of Activities")
    rows = []
    rows.append(("REVENUE AND SUPPORT", ""))
    for label, val in d["revenue"].items():
        rows.append((f"    {label}", _fmt(val)))
    rows.append(("Total Revenue and Support", _fmt(d["total_revenue"])))
    rows.append(("", ""))
    rows.append(("EXPENSES", ""))
    for label, val in d["expenses"].items():
        rows.append((f"    {label}", _fmt(val)))
    rows.append(("Total Expenses", _fmt(d["total_expenses"])))
    rows.append(("", ""))
    rows.append(("CHANGE IN NET ASSETS", _fmt(d["change_in_net_assets"])))
    _add_table_slide(prs, f"Statement of Activities — {org_name}",
                     ["Description", "Amount"], rows, col_widths=[6.5, 2.5])

    # --- Statement of Financial Position ---
    _add_section_slide(prs, "Statement of Financial Position")
    rows = []
    rows.append(("ASSETS", ""))
    for label, val in pos["assets"].items():
        prefix = "" if "Total" in label else "    "
        rows.append((f"{prefix}{label}", _fmt(val)))
    rows.append(("", ""))
    rows.append(("LIABILITIES", ""))
    for label, val in pos["liabilities"].items():
        rows.append((f"    {label}", _fmt(val)))
    rows.append(("", ""))
    rows.append(("NET ASSETS", ""))
    for label, val in pos["net_assets"].items():
        prefix = "" if "Total" in label else "    "
        rows.append((f"{prefix}{label}", _fmt(val)))
    rows.append(("", ""))
    rows.append(("TOTAL LIABILITIES AND NET ASSETS",
                 _fmt(pos["total_liabilities_and_net_assets"])))
    _add_table_slide(prs, f"Statement of Financial Position — {org_name}",
                     ["Description", "Amount"], rows, col_widths=[6.5, 2.5])

    # --- Statement of Functional Expenses ---
    if func["table"]:
        _add_section_slide(prs, "Statement of Functional Expenses")
        headers = ["Category"] + func["functional_categories"] + ["Total"]
        rows = []
        for nat_cat, values in func["table"].items():
            row = [nat_cat]
            for fc in func["functional_categories"]:
                row.append(_fmt(values.get(fc, 0)))
            row.append(_fmt(values.get("Total", 0)))
            rows.append(tuple(row))
        total_row = ["TOTAL EXPENSES"]
        for fc in func["functional_categories"]:
            total_row.append(_fmt(func["totals"].get(fc, 0)))
        total_row.append(_fmt(func["totals"].get("Total", 0)))
        rows.append(tuple(total_row))
        n_cols = len(headers)
        cw = [3.5] + [1.5] * (n_cols - 1)
        _add_table_slide(prs, f"Functional Expenses — {org_name}",
                         headers, rows, col_widths=cw)

    # --- Statement of Cash Flows ---
    _add_section_slide(prs, "Statement of Cash Flows")
    rows = []
    rows.append(("OPERATING ACTIVITIES", ""))
    for label, val in cf["operating_activities"]["inflows"].items():
        rows.append((f"    {label}", _fmt(val)))
    for label, val in cf["operating_activities"]["outflows"].items():
        rows.append((f"    {label}", _fmt(val)))
    rows.append(("Net Cash from Operating Activities",
                 _fmt(cf["operating_activities"]["net"])))
    rows.append(("", ""))
    rows.append(("Net Cash from Investing Activities",
                 _fmt(cf["investing_activities"]["net"])))
    rows.append(("Net Cash from Financing Activities",
                 _fmt(cf["financing_activities"]["net"])))
    rows.append(("", ""))
    rows.append(("Net Change in Cash", _fmt(cf["net_change_in_cash"])))
    rows.append(("Beginning Cash Balance", _fmt(cf["beginning_cash"])))
    rows.append(("ENDING CASH BALANCE", _fmt(cf["ending_cash"])))
    _add_table_slide(prs, f"Statement of Cash Flows — {org_name}",
                     ["Description", "Amount"], rows, col_widths=[6.5, 2.5])

    # --- Notes / blank slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    p = txBox.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Notes"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = _DARK
    run.font.name = "Calibri"

    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.4))
    p2 = txBox2.text_frame.paragraphs[0]
    r2 = p2.add_run()
    r2.text = "Use this space for additional notes, charts, or supporting documentation."
    r2.font.size = Pt(11)
    r2.font.italic = True
    r2.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    r2.font.name = "Calibri"

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()


# ---------------------------------------------------------------------------
# PDF export
# ---------------------------------------------------------------------------

class _FinancialPDF(FPDF):
    def __init__(self, org_name="Organization"):
        super().__init__(orientation="P", unit="mm", format="letter")
        self.org_name = org_name
        self.set_auto_page_break(auto=True, margin=20)

    def header(self):
        if self.page_no() == 1:
            return
        self.set_font("Helvetica", "I", 8)
        self.set_text_color(120, 120, 120)
        self.cell(0, 5, self.org_name, align="L")
        self.cell(0, 5, f"Page {self.page_no()}", align="R", new_x="LMARGIN", new_y="NEXT")
        self.line(10, 12, self.w - 10, 12)
        self.ln(3)

    def footer(self):
        self.set_y(-15)
        self.set_font("Helvetica", "I", 7)
        self.set_text_color(150, 150, 150)
        self.cell(0, 10, "Generated by Nonprofit Financial Analyzer", align="C")

    def add_cover(self, title, period):
        self.add_page()
        self.ln(60)
        self.set_font("Helvetica", "B", 28)
        self.set_text_color(45, 50, 86)
        self.cell(0, 15, self.org_name, align="C", new_x="LMARGIN", new_y="NEXT")
        self.ln(5)
        self.set_font("Helvetica", "", 16)
        self.set_text_color(100, 100, 140)
        self.cell(0, 10, title, align="C", new_x="LMARGIN", new_y="NEXT")
        self.ln(3)
        self.set_font("Helvetica", "I", 12)
        self.cell(0, 8, period, align="C", new_x="LMARGIN", new_y="NEXT")

    def add_statement_header(self, title, subtitle=""):
        self.add_page()
        self.set_font("Helvetica", "B", 16)
        self.set_text_color(45, 50, 86)
        self.cell(0, 10, title, align="C", new_x="LMARGIN", new_y="NEXT")
        if subtitle:
            self.set_font("Helvetica", "I", 10)
            self.set_text_color(100, 100, 100)
            self.cell(0, 6, subtitle, align="C", new_x="LMARGIN", new_y="NEXT")
        self.ln(2)
        self.line(10, self.get_y(), self.w - 10, self.get_y())
        self.ln(4)

    def add_section_label(self, text):
        self.set_font("Helvetica", "B", 10)
        self.set_text_color(50, 50, 50)
        self.set_fill_color(235, 237, 248)
        self.cell(0, 7, text, new_x="LMARGIN", new_y="NEXT", fill=True)
        self.ln(1)

    def add_line_item(self, label, amount, bold=False, indent=0):
        style = "B" if bold else ""
        self.set_font("Helvetica", style, 10)
        self.set_text_color(50, 50, 50)
        x = 12 + indent * 6
        self.set_x(x)
        label_w = 130 - indent * 6
        self.cell(label_w, 6, label)
        self.set_font("Helvetica", style, 10)
        self.cell(50, 6, _fmt(amount), align="R", new_x="LMARGIN", new_y="NEXT")

    def add_total_line(self, label, amount):
        y = self.get_y()
        self.line(10, y, self.w - 10, y)
        self.ln(1)
        self.add_line_item(label, amount, bold=True)
        y2 = self.get_y()
        self.line(10, y2, self.w - 10, y2)
        self.ln(2)

    def add_spacer(self, h=3):
        self.ln(h)

    def add_chart_image(self, img_bytes, title=""):
        if title:
            self.set_font("Helvetica", "B", 12)
            self.set_text_color(45, 50, 86)
            self.cell(0, 8, title, align="C", new_x="LMARGIN", new_y="NEXT")
            self.ln(2)
        img_stream = io.BytesIO(img_bytes)
        page_w = self.w - 20
        self.image(img_stream, x=10, w=page_w)
        self.ln(5)


def export_to_pdf(statements: dict, transactions_df: pd.DataFrame,
                  org_name: str = "Organization") -> bytes:
    """Export financial statements and charts to PDF."""
    act = statements["activities"]
    pos = statements["position"]
    func = statements["functional_expenses"]
    cf = statements["cash_flows"]
    d = act["without_donor_restrictions"]
    period = act.get("period", "")

    pdf = _FinancialPDF(org_name)

    # --- Cover page ---
    pdf.add_cover("Financial Statements", period)

    # --- Summary metrics ---
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 16)
    pdf.set_text_color(45, 50, 86)
    pdf.cell(0, 10, "Financial Summary", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(6)

    metrics = [
        ("Total Revenue", d["total_revenue"]),
        ("Total Expenses", d["total_expenses"]),
        ("Change in Net Assets", d["change_in_net_assets"]),
        ("Ending Cash", cf["ending_cash"]),
    ]
    col_w = (pdf.w - 20) / len(metrics)
    x_start = 10
    for i, (label, value) in enumerate(metrics):
        x = x_start + i * col_w
        pdf.set_fill_color(45, 50, 86)
        pdf.set_xy(x, pdf.get_y())
        pdf.rect(x, pdf.get_y(), col_w - 2, 20, style="F")
        pdf.set_xy(x, pdf.get_y())
        pdf.set_font("Helvetica", "", 8)
        pdf.set_text_color(200, 200, 220)
        pdf.cell(col_w - 2, 8, label, align="C")
        pdf.set_xy(x, pdf.get_y() + 8)
        pdf.set_font("Helvetica", "B", 13)
        pdf.set_text_color(255, 255, 255)
        pdf.cell(col_w - 2, 12, _fmt(value), align="C")

    pdf.set_y(pdf.get_y() + 28)
    pdf.set_text_color(50, 50, 50)

    # --- Charts ---
    chart_builders = [
        ("Revenue by Category", _build_revenue_pie),
        ("Expenses by Category", _build_expense_pie),
        ("Monthly Cash Flow", _build_monthly_bar),
    ]
    for chart_title, builder in chart_builders:
        fig = builder(transactions_df)
        if fig is not None:
            fig.update_layout(width=1400, height=600)
            img_bytes = _try_chart_image(fig, width=1400, height=600)
            if img_bytes:
                pdf.add_page()
                pdf.add_chart_image(img_bytes, chart_title)

    # --- Statement of Activities ---
    pdf.add_statement_header("Statement of Activities", period)

    pdf.add_section_label("REVENUE AND SUPPORT")
    for label, val in d["revenue"].items():
        pdf.add_line_item(label, val, indent=1)
    pdf.add_total_line("Total Revenue and Support", d["total_revenue"])

    pdf.add_section_label("EXPENSES")
    for label, val in d["expenses"].items():
        pdf.add_line_item(label, val, indent=1)
    pdf.add_total_line("Total Expenses", d["total_expenses"])

    if d.get("net_transfers", 0) != 0:
        pdf.add_section_label("INTERNAL ACCOUNT TRANSFERS (NET)")
        pdf.add_line_item("Net Transfers", d["net_transfers"], indent=1)
        pdf.add_spacer()

    pdf.add_spacer(5)
    pdf.add_line_item("CHANGE IN NET ASSETS", d["change_in_net_assets"], bold=True)

    # --- Statement of Financial Position ---
    pdf.add_statement_header("Statement of Financial Position",
                             pos.get("as_of", ""))

    pdf.add_section_label("ASSETS")
    for label, val in pos["assets"].items():
        is_total = "Total" in label
        pdf.add_line_item(label, val, bold=is_total, indent=0 if is_total else 1)
    pdf.add_spacer()

    pdf.add_section_label("LIABILITIES")
    for label, val in pos["liabilities"].items():
        pdf.add_line_item(label, val, indent=1)
    pdf.add_spacer()

    pdf.add_section_label("NET ASSETS")
    for label, val in pos["net_assets"].items():
        is_total = "Total" in label
        pdf.add_line_item(label, val, bold=is_total, indent=0 if is_total else 1)
    pdf.add_spacer()
    pdf.add_total_line("TOTAL LIABILITIES AND NET ASSETS",
                       pos["total_liabilities_and_net_assets"])

    # --- Statement of Functional Expenses ---
    if func["table"]:
        pdf.add_statement_header("Statement of Functional Expenses", period)

        func_cats = func["functional_categories"]
        all_cols = ["Category"] + func_cats + ["Total"]
        n_cols = len(all_cols)
        cat_w = 50
        num_w = (pdf.w - 20 - cat_w) / (n_cols - 1)

        pdf.set_font("Helvetica", "B", 8)
        pdf.set_fill_color(68, 82, 122)
        pdf.set_text_color(255, 255, 255)
        pdf.cell(cat_w, 7, "Category", border=1, fill=True)
        for col in func_cats + ["Total"]:
            short = col.replace("Management & General", "Mgmt & Gen")
            pdf.cell(num_w, 7, short, border=1, fill=True, align="C")
        pdf.ln()

        alt = False
        for nat_cat, values in func["table"].items():
            if alt:
                pdf.set_fill_color(240, 240, 248)
            else:
                pdf.set_fill_color(255, 255, 255)
            pdf.set_text_color(50, 50, 50)
            pdf.set_font("Helvetica", "", 8)
            pdf.cell(cat_w, 6, nat_cat, border=1, fill=True)
            for fc in func_cats + ["Total"]:
                pdf.cell(num_w, 6, _fmt(values.get(fc, 0)), border=1,
                         fill=True, align="R")
            pdf.ln()
            alt = not alt

        pdf.set_font("Helvetica", "B", 8)
        pdf.set_fill_color(232, 234, 246)
        pdf.cell(cat_w, 7, "TOTAL EXPENSES", border=1, fill=True)
        for fc in func_cats + ["Total"]:
            pdf.cell(num_w, 7, _fmt(func["totals"].get(fc, 0)), border=1,
                     fill=True, align="R")
        pdf.ln()

    # --- Statement of Cash Flows ---
    pdf.add_statement_header("Statement of Cash Flows", period)

    pdf.add_section_label("CASH FLOWS FROM OPERATING ACTIVITIES")
    for label, val in cf["operating_activities"]["inflows"].items():
        pdf.add_line_item(label, val, indent=1)
    for label, val in cf["operating_activities"]["outflows"].items():
        pdf.add_line_item(label, val, indent=1)
    pdf.add_line_item("Net Cash from Operating Activities",
                      cf["operating_activities"]["net"], bold=True)
    pdf.add_spacer()

    pdf.add_section_label("CASH FLOWS FROM INVESTING ACTIVITIES")
    pdf.add_line_item("Net Cash from Investing Activities",
                      cf["investing_activities"]["net"], bold=True)
    pdf.add_spacer()

    pdf.add_section_label("CASH FLOWS FROM FINANCING ACTIVITIES")
    pdf.add_line_item("Net Cash from Financing Activities",
                      cf["financing_activities"]["net"], bold=True)
    pdf.add_spacer(5)

    pdf.add_line_item("Net Change in Cash", cf["net_change_in_cash"], bold=True)
    pdf.add_line_item("Beginning Cash Balance", cf["beginning_cash"])
    pdf.add_total_line("ENDING CASH BALANCE", cf["ending_cash"])

    # --- Notes / blank page ---
    pdf.add_statement_header("Notes")
    pdf.set_font("Helvetica", "I", 10)
    pdf.set_text_color(150, 150, 150)
    pdf.cell(0, 8, "Use this space for additional notes, charts, or supporting documentation.",
             new_x="LMARGIN", new_y="NEXT")

    output = io.BytesIO()
    pdf.output(output)
    output.seek(0)
    return output.getvalue()
